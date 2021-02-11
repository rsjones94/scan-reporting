#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Created on Thu Jul 30 12:38:31 2020

@author: manusdonahue
"""

import os
import time
from PIL import Image
import shutil
import itertools
import operator
import re

import numpy as np
from pptx import Presentation
from pptx.util import Inches
import pandas as pd


def replace_in_ppt(search_str, repl_str, filename):
    """"search and replace text in PowerPoint while preserving formatting"""
    #Useful Links ;)
    #https://stackoverflow.com/questions/37924808/python-pptx-power-point-find-and-replace-text-ctrl-h
    #https://stackoverflow.com/questions/45247042/how-to-keep-original-text-formatting-of-text-with-python-powerpoint
    prs = Presentation(filename)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if(shape.text.find(search_str))!=-1:
                    text_frame = shape.text_frame
                    cur_text = text_frame.paragraphs[0].runs[0].text
                    new_text = cur_text.replace(str(search_str), str(repl_str))
                    text_frame.paragraphs[0].runs[0].text = new_text
    prs.save(filename)


def analyze_ppt(inp, output):
    """ Take the input file and analyze the structure.
    The output file contains marked up information to make it easier
    for generating future powerpoint templates.
    """
    prs = Presentation(inp)
    # Each powerpoint file has multiple layouts
    # Loop through them all and  see where the various elements are
    for index, _ in enumerate(prs.slide_layouts):
        slide = prs.slides.add_slide(prs.slide_layouts[index])
        # Not every slide has to have a title
        try:
            title = slide.shapes.title
            title.text = 'Title for Layout {}'.format(index)
        except AttributeError:
            print("No Title for Layout {}".format(index))
        # Go through all the placeholders and identify them by index and type
        for shape in slide.placeholders:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                # Do not overwrite the title which is just a special placeholder
                try:
                    if 'Title' not in shape.text:
                        shape.text = 'Placeholder index:{} type:{}'.format(phf.idx, shape.name)
                except AttributeError:
                    print("{} has no text attribute".format(phf.type))
                print('{} {}'.format(phf.idx, shape.name))
    prs.save(output)


def plot_dot(slide, dot_img, x, y, origin, xpi, ypi, size=0.1):
    """
    A very specific function for putting a dot on an image of a line plot
    in a powerpoint.
    

    Parameters
    ----------
    slide : pptx slide object
        the slide you're plotting in.
    dot_img : str
        path to image you want to place
    x : float
        the x value to plot.
    y : float
        the y value to plot.
    origin : tuple of int
        x,y of the origin of the chart in inches.
    xpi : float
        the number of x units per inch.
    ypi : float
        the number of y units per inch.
    size: float
        size of the inserted dot image in inches

    Returns
    -------
    None.

    """
    
    # images are placed using the coords of their upper left corner
    
    ox, oy = origin
    
    off_x = x / xpi # raw offset from origin in inches
    off_y = y / ypi
    
    over_x = ox + off_x # absolute position after accounting for origin but not image size
    over_y = oy - off_y
    
    account_x = over_x - (size/2)
    account_y = over_y - (size/2)
    
    plot_x = account_x
    plot_y = account_y

    shp = slide.shapes
    #print(f'Plotting {plot_x,plot_y}. Origin: {origin}. xy: {x,y}')
    picture = shp.add_picture(dot_img, Inches(plot_x), Inches(plot_y), width=Inches(size), height=Inches(size))


def add_ppt_image(slide, img, scale=0.3, insert_type='img', poster=None, at=(0,0)):
    ex, why = at
    shp = slide.shapes
    
    if insert_type=='img':
        im = Image.open(img)
        width, height = im.size
        picture = shp.add_picture(img, Inches(ex), Inches(why))
        picture.width = int(picture.width*scale)
        picture.height = int(picture.height*scale)
    elif insert_type=='mov':
        picture = slide.shapes.add_movie(img, ex, why, 20000, 20000, poster_frame_image=poster, mime_type='video/mp4')
        

def add_ppt_image_ph(slide, placeholder_id, image_url):
    placeholder = slide.placeholders[placeholder_id]
 
    # Calculate the image size of the image
    im = Image.open(image_url)
    width, height = im.size
 
    # Make sure the placeholder doesn't zoom in
    placeholder.height = height
    placeholder.width = width
 
    # Insert the picture
    placeholder = placeholder.insert_picture(image_url)
    
    # Calculate ratios and compare
    image_ratio = width / height
    placeholder_ratio = placeholder.width / placeholder.height
    ratio_difference = placeholder_ratio - image_ratio
 
    # Placeholder width too wide:
    if ratio_difference > 0:
        difference_on_each_side = ratio_difference / 2
        placeholder.crop_left = -difference_on_each_side
        placeholder.crop_right = -difference_on_each_side
    # Placeholder height too high
    else:
        difference_on_each_side = -ratio_difference / 2
        placeholder.crop_bottom = -difference_on_each_side
        placeholder.crop_top = -difference_on_each_side


def get_terminal(path):
    """
    Takes a filepath or directory tree and returns the last file or directory
    

    Parameters
    ----------
    path : path
        path in question.

    Returns
    -------
    str of only the final file or directory.

    """
    
    return os.path.basename(os.path.normpath(path))


def str_time_elapsed(start):
    now = time.time()
    out = now - start
    pretty = round(out/60, 2)
    return pretty


def any_in_str(s, l):
    """
    Returns whether any of a list of substrings is in a string
    

    Parameters
    ----------
    s : str
        string to look for substrings in.
    l : list of str
        substrings to check for in s.

    Returns
    -------
    bool

    """
    
    return any([substr in s for substr in l])


def dicom_to_parrec(filename, out_folder, path_to_perl_script='/Users/manusdonahue/Desktop/Projects/gstudy_converter/convert_dicom_to_xmlrec.pl'):
    """
    Calls Brian Welch's perl script that converts DICOMs to PARRECs (using default flags).
    
    The perl script by default creates 4 files (PAR, REC, V41 and XML) in a subfolder
    within the directory that the target file is called xmlparrec, and the original
    DICOM is renamed in the process. THIS SCRIPT creates a copy of the original DICOM
    and destroys it and the xmlparrec folder after moving the output to out_folder.
    
    Note that this script uses the default renaming convention:
        $Patient_Name$_$%02d%Acquisition_Number$_$%02d%Reconstruction_Number$_$SeriesTime$_($Protocol_Name$)
    
    Also note that you probably can't anticipate what the output names will be,
    since that would require programatically inspecting the DICOM metedata. If you
    could/wanted to do this you probably wouldn't be using this wrapper.
    

    Parameters
    ----------
    filename : pathlike
        path to the dicom you wish to convert.
    out_folder : pathlike
        the directory to which the PARREC will be written.
    path_to_perl_script : pathlike
        path to the gstudy converter perl script.

    Returns
    -------
    None.

    """
    new_wd = os.path.dirname(os.path.normpath(path_to_perl_script))
    
    old_wd = os.getcwd()
    os.chdir(new_wd)
    
    parent_folder = os.path.dirname(os.path.normpath(filename))
    file_basename = os.path.basename(os.path.normpath(filename))
    tmp_folder = os.path.join(parent_folder, 'tmp')
    xml_folder = os.path.join(tmp_folder, 'xmlparrec')
    
    if os.path.exists(tmp_folder):
        shutil.rmtree(tmp_folder)
    os.mkdir(tmp_folder)
    copyname = os.path.join(tmp_folder, file_basename)
    
    shutil.copyfile(filename, copyname)
    
    try:
    
        call = f'perl {path_to_perl_script} -d {tmp_folder} -f {copyname}'
        #print(f'Call: {call}')
        os.system(call)
        
        conv_files = [f for f in os.listdir(xml_folder) if os.path.isfile(os.path.join(xml_folder, f))]
        assert len(conv_files) == 4
        
        for fi in conv_files:
            full_name = os.path.join(xml_folder, fi)
            target = os.path.join(out_folder, fi)
            shutil.move(full_name, target)
            
    except AssertionError:
        print(f'Expected 4 files in {xml_folder}, but found {len(conv_files)}. Cleaning up....')
    
    finally:
        shutil.rmtree(tmp_folder)
        os.chdir(old_wd)
    
    
def most_common(L):
    """
    Courtesy Alex Martelli
    """
    # get an iterable of (item, iterable) pairs
    SL = sorted((x, i) for i, x in enumerate(L))
    # print 'SL:', SL
    groups = itertools.groupby(SL, key=operator.itemgetter(0))
    # auxiliary function to get "quality" for an item
    def _auxfun(g):
      item, iterable = g
      count = 0
      min_index = len(L)
      for _, where in iterable:
        count += 1
        min_index = min(min_index, where)
      # print 'item %r, count %r, minind %r' % (item, count, min_index)
      return count, -min_index
    # pick the highest-count/earliest item
    return max(groups, key=_auxfun)[0]


def parse_scd_csv(in_csv, scan_index, std=False):
    """
    
    

    Parameters
    ----------
    in_csv : TYPE
        DESCRIPTION.
    scan_index : TYPE
        DESCRIPTION.
    std : BOOL, optional
        If False, then _cbf metrics are CBF. If True, then they are standard deviations. The default is False.

    Returns
    -------
    parsed : TYPE
        DESCRIPTION.

    """
    
    if std:
        row_add = 3
        
    else:
        row_add = 0
    
    raw = open(in_csv).readlines()
    
    stripped_to_nums = [re.sub('[^0-9,.-]', "", i).split(',') for i in raw] # regex to leave only numbers, commas, periods and dashes, then split by commas
    
    for i, big_list in enumerate(stripped_to_nums):
        for j, sub in enumerate(big_list):
            if sub == '-999':
                stripped_to_nums[i][j] = np.nan
    
    parsed = {
                    f'mr{scan_index+1}_lparietal_gm_cbf':str(stripped_to_nums[15+row_add][0]),
                    f'mr{scan_index+1}_rparietal_gm_cbf':str(stripped_to_nums[15+row_add][1]),
                    f'mr{scan_index+1}_lfrontal_gm_cbf':str(stripped_to_nums[15+row_add][2]),
                    f'mr{scan_index+1}_rfrontal_gm_cbf':str(stripped_to_nums[15+row_add][3]),
                    f'mr{scan_index+1}_loccipital_gm_cbf':str(stripped_to_nums[15+row_add][4]),
                    f'mr{scan_index+1}_roccipital_gm_cbf':str(stripped_to_nums[15+row_add][5]),
                    f'mr{scan_index+1}_ltemporal_gm_cbf':str(stripped_to_nums[15+row_add][6]),
                    f'mr{scan_index+1}_rtemporal_gm_cbf':str(stripped_to_nums[15+row_add][7]),
                    f'mr{scan_index+1}_lcerebellum_gm_cbf':str(stripped_to_nums[15+row_add][8]),
                    f'mr{scan_index+1}_rcerebellum_gm_cbf':str(stripped_to_nums[15+row_add][9]),
                    f'mr{scan_index+1}_recalc_gm_cbf':str(stripped_to_nums[15+row_add][10]),
                    f'mr{scan_index+1}_recalc_wm_cbf':str(stripped_to_nums[15+row_add][11]),
                    f'mr{scan_index+1}_white_cbv':str(stripped_to_nums[7][0]),
                    f'mr{scan_index+1}_grey_cbv':str(stripped_to_nums[8][0]),
                    f'mr{scan_index+1}_csf_cbv':str(stripped_to_nums[9][0]),
                    f'mr{scan_index+1}_relaxation_rate1':str(1/float(stripped_to_nums[2][1])),
                    f'mr{scan_index+1}_relaxation_rate2':str(1/float(stripped_to_nums[2][3])),
                    f'mr{scan_index+1}_venous_oxygen_sat1':str(float(stripped_to_nums[3][1])*100), #bovine
                    f'mr{scan_index+1}_venous_oxygen_sat2':str(float(stripped_to_nums[3][3])*100), #bovine
                    f'mr{scan_index+1}_aa_model_venous_oxygen_sat1':str(float(stripped_to_nums[5][1])*100),
                    f'mr{scan_index+1}_aa_model_venous_oxygen_sat2':str(float(stripped_to_nums[5][3])*100),
                    f'mr{scan_index+1}_ss_model_venous_oxygen_sat1':str(float(stripped_to_nums[6][1])*100),
                    f'mr{scan_index+1}_ss_model_venous_oxygen_sat2':str(float(stripped_to_nums[6][3])*100),
                    f'mr{scan_index+1}_f_model_venous_oxygen_sat1':str(float(stripped_to_nums[4][1])*100),
                    f'mr{scan_index+1}_f_model_venous_oxygen_sat2':str(float(stripped_to_nums[4][3])*100)
                    
                 }
    
    return parsed
    