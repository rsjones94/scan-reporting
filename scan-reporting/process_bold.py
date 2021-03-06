#!/usr/bin/env python3

# -*- coding: utf-8 -*-
help_info = """
This script takes raw scan data for BOLD scans and processes it completely.
This involves:
    1) deidentifying the scans
    2) generating derived images - CBF, CVR, CVR_Max, and CVR_Delay images in native space, T1 space, and 4 mm MNI space.
    3) generating the CVR movie
    4) metrics calculation
    5) reporting image generation
    6) creating patient scan report as a powerpoint
    
Note that to have the EtCO2 trace generated for you (step 6), you should add
a CSV with two columns in the PTSTEN folder (not the Acquired subdirectory)
called ‘etco2.csv’. Each column should have a header (though the header name
does not matter). The first column should be the dynamic scan numbers,
and the second column should be the EtCO2 values.
    
    
input:
    -i / --infolder : the path to the folder of pt data. of form /Users/manusdonahue/Desktop/Projects/BOLD/Data/[PTSTEN_ID]
        this folder should have one subfolder, Acquired, which contains the raw scans
    -n / --name : the name of the patient to deidentify. Only required if running step 1
    -s / --steps : the steps to actually calculate, passed as a string of numbers
        if 0, all steps will be carried out. Otherwise, only the steps that appear
        in the string will be carried out. For example, 134 will cause the script
        to only go through steps 1, 3 and 4
    -d / --dob : the date of birth as YYYY.mm.dd, or their age as a float
        Only used in step 6. If -d is a date, the dob, age and metric plotting
        is completed in the template for you (age is calculated as differnce between
        scan date in PAR files and DOB). If -d is an age, only the age and metric
        plotting is completed. If -d is not supplied, then neither the dob or age is completed,
        and plotting is carried out using age=0.
    -g / --help : brings up this helpful information. does not take an argument
"""

import os
import sys
import getopt
import subprocess
import time
import datetime
import glob
import shutil
import datetime

from pptx import Presentation
import matplotlib
import matplotlib.pyplot as plt
import pandas as pd
import nibabel as nib

from helpers import get_terminal, str_time_elapsed, any_in_str, replace_in_ppt, analyze_ppt, add_ppt_image, add_ppt_image_ph, plot_dot
import helpers as hp
from report_image_generation import par2nii, nii_image

#sys.exit()

inp = sys.argv
bash_input = inp[1:]
options, remainder = getopt.getopt(bash_input, "i:n:s:d:c:g", ["infolder=","name=",'steps=','dob=', 'clean', 'help'])

for opt, arg in options:
    if opt in ('-i', '--infile'):
        in_folder = arg
    elif opt in ('-n', '--name'):
        deidentify_name = arg
    elif opt in ('-s', '--steps'):
        steps = arg
    elif opt in ('-d', '--dob'):
        dobage = arg
    elif opt in ('-g', '--help'):
        print(help_info)
        sys.exit()

try:
    if steps == '0':
        steps = '123456'
except NameError:
    print('-s not specified. running all steps')
    steps = '123456'
        
try:
    assert os.path.isdir(in_folder)
except AssertionError:
    raise AssertionError('input folder does not exist')
    
    




start_stamp = time.time()
now = datetime.datetime.now()
pretty_now = now.strftime("%Y-%m-%d %H:%M:%S")

inp_copy = inp.copy()
for i, s in enumerate(inp_copy):
    if s == '-n' or s== '--name':
        inp_copy[i+1] = '[REDACTED]'

thecommand = ' '.join(inp_copy)
meta_file_name = os.path.join(in_folder, 'meta.txt')
meta_file = open(meta_file_name, 'w')
meta_file.write(f'Processing started {pretty_now}')
meta_file.write('\n\n')
meta_file.write(thecommand)
meta_file.close()

print(f'\nBegin processing: {pretty_now}')






acq_folder = os.path.join(in_folder, 'Acquired')
orig_files = [os.path.join(acq_folder, f) for f in os.listdir(acq_folder) if os.path.isfile(os.path.join(acq_folder, f))]
extensions = [f.split('.')[-1] for f in orig_files]
guess_ext = hp.most_common(extensions)

orig_data_copy_folder = os.path.join(in_folder, 'rawdata')
    
dcm_exts = ['dcm', 'DCM']
parrec_exts = ['PAR', 'REC', 'V41', 'XML']
nii_exts = ['nii', 'gz']

if guess_ext in parrec_exts:
    print('Input files seem to be PARREC - proceeding as normal')
elif guess_ext in dcm_exts:
    print('Input files seem to be DICOM - converting to PARREC before continuing (original DICOMs will be retained)')
    shutil.copytree(acq_folder, orig_data_copy_folder)
    shutil.rmtree(acq_folder)
    os.mkdir(acq_folder)
    
    moved_files = [os.path.join(orig_data_copy_folder, f) for f in os.listdir(orig_data_copy_folder) if os.path.isfile(os.path.join(orig_data_copy_folder, f))]
    moved_extensions = [f.split('.')[-1] for f in orig_files]
    for fi, ext in zip(moved_files, moved_extensions):
        if ext in dcm_exts:
            #print(f'\n\n\nCONVERTING: {fi}')
            hp.dicom_to_parrec(fi, acq_folder)
elif guess_ext in nii_exts:
    has_ans = False
    while not has_ans:
        ans = input(f'Input files seem to be NiFTI. ASL processing of NiFTIs is in an UNSTABLE BETA state.\nRESULTS MUST BE MANUALLY INSPECTED FOR CORRECTNESS. Please acknowledge this or cancel processing. [acknowledge/cancel]\n')
        if ans in ('acknowledge', 'cancel'):
            has_ans = True
            if ans == 'cancel':
                raise Exception('Aborting processing')
            elif ans == 'acknowledge':
                print('Continuing with beta processing of NiFTIs')
        else:
            print('Answer must be "acknowledge" or "cancel"')
else:
    raise Exception(f'Filetype ({guess_ext}) does not seem to be supported')

if guess_ext in nii_exts:
    fig_ext = 'nii.gz'
else:
    fig_ext = guess_ext


original_wd = os.getcwd()
pt_id = replacement = get_terminal(in_folder) # if the input folder is named correctly, it is the ID that will replace the pt name

if '1' in steps:
    ##### step 1 : deidentification
    
    assert type(deidentify_name) == str, 'patient name must be a string'
        
    files_of_interest = os.listdir(os.path.join(in_folder, 'Acquired'))
    has_deid_name = any([deidentify_name in f for f in files_of_interest])
    if not has_deid_name:
        has_ans = False
        while not has_ans:
            ans = input(f'\nName "{deidentify_name}" not found in Acquired folder. Would you like to proceed anyway? [y/n/change]\n')
            if ans in ('y','n', 'change'):
                has_ans = True
                if ans == 'n':
                    raise Exception('Aborting processing')
                elif ans == 'change':
                    deidentify_name = input(f'Enter a new deidentification string to replace {deidentify_name}:\n')
            else:
                print('Answer must be "y", "n" or "change')
        
    print(f'\nStep 1: deidentification. {deidentify_name} will be replaced with {replacement}')
        
    # build the call to the deidentify script
    deid_scripts_loc = r'/Users/manusdonahue/Desktop/Projects/BOLD/Scripts/deidentifySLW'
    
    strip_filename_input = f'deidentifyFileNames.sh {in_folder} {deidentify_name} {replacement}'
    strip_header_input = f'deidentifyPARfiles.sh {in_folder}'
    
    deid_commands = [os.path.join(deid_scripts_loc, c) for c in (strip_filename_input, strip_header_input)]
    
    for com in deid_commands:
        subprocess.run([com], check=True, shell=True)
        
    print(f'\nDeidentification complete. Elapsed time: {str_time_elapsed(start_stamp)} minutes')

if '2' in steps:    
    ##### step 2 : main processing
    
    print(f'Step 2: begin main processing sequence\n')
    
    asltype = 'Baseline'
    dynamics = 360
    

    has_ans = False
    while not has_ans:
        ans = input('What is the ASL type? [PASL / pCASL]\n')
        if ans in ('pCASL', 'PASL'):
            has_ans = True
            confirm = input(f'Please confirm that ASL type is {ans} by entering the ASL type again\n')
            if ans != confirm:
                has_ans = False
                print(f'\nConfirmation failed ({ans} != {confirm})\n')
            else:
                print('\nEntry confirmed\n')
        else:
            print('Answer must be PASL or pCASL')
            
    if ans == 'pCASL':
        pcaslBool = 1
    elif ans == 'PASL':
        pcaslBool = 0
    
    
    processing_scripts_loc = r'/Users/manusdonahue/Desktop/Projects/BOLD/Scripts/'
    os.chdir(processing_scripts_loc)
    processing_input = f'''/Applications/MATLAB_R2016b.app/bin/matlab -nodesktop -nosplash -r "Master('{pt_id}','{asltype}',{dynamics},{pcaslBool})"'''
    
    print(f'Call to MATLAB: {processing_input}')
    
    subprocess.run(processing_input, check=True, shell=True)
    # subprocess.run('quit force', check=True, shell=True)
    
    os.chdir(original_wd)
    
    print(f'\nMain processing complete. Elapsed time: {str_time_elapsed(start_stamp)} minutes')
    
if '3' in steps:
    ##### step 3 : generate cvr movie
    
    print(f'\nStep 3: generating CVR movie')
    
    movie_scripts_loc = r'/Users/manusdonahue/Desktop/Projects/BOLD/Scripts/zstatMov'
    os.chdir(movie_scripts_loc)
    
    pt1 = f'./mkZstatMov_part1_v2.sh /Users/manusdonahue/Desktop/Projects/BOLD/Data/{pt_id}'
    subprocess.run(pt1, check=True, shell=True)
    
    mov_filepath = f'/Users/manusdonahue/Desktop/Projects/BOLD/Data/{pt_id}/{pt_id}_zstatMovie.nii.gz'
    
    pt2 = f'''/Applications/MATLAB_R2016b.app/bin/matlab -nodesktop -nosplash -r "mkZstatMov_part2_v4('{mov_filepath}')"'''
    subprocess.run(pt2, check=True, shell=True)
    # subprocess.run('quit force', check=True, shell=True)
    
    os.chdir(original_wd)
    
    print(f'\nCVR movie generation complete. Elapsed time: {str_time_elapsed(start_stamp)} minutes')
    

if '4' in steps:
    ##### step 4 : metrics calculation
    
    print(f'\nStep 4: calculating metrics')
    
    processing_scripts_loc = r'/Users/manusdonahue/Desktop/Projects/BOLD/Scripts/'
    os.chdir(processing_scripts_loc)
    
    cmd = f'''/Applications/MATLAB_R2016b.app/bin/matlab -nodesktop -nosplash -r "Calculate_Metrics('{pt_id}', 1)"'''
    subprocess.run(cmd, check=True, shell=True)
    
    os.chdir(original_wd)
    
    print(f'\nMetrics calculation complete. Elapsed time: {str_time_elapsed(start_stamp)} minutes')
    

signature_relationships = {('FLAIR_AX', 'T2W_FLAIR'):
                               {'basename': 'axFLAIR', 'excl':['cor','COR','coronal','CORONAL'], 'isin':'Acquired', 'ext':fig_ext, 'cmap':matplotlib.cm.gray, 'dims':(4,6)}, # THIS NEEDS TO BE UPDATED - the input FLAIR will not always be PAR!
                           ('CBF_MNI',):
                               {'basename': 'CBF', 'excl':[], 'isin':'processed', 'ext':'nii.gz', 'cmap':matplotlib.cm.jet, 'dims':(3,10)},
                           ('ZSTAT1_MNI_normalized',):
                               {'basename': 'CVR', 'excl':[], 'isin':'processed', 'ext':'nii.gz', 'cmap':matplotlib.cm.jet, 'dims':(3,10)},
                           ('ZMAX2STANDARD_normalized',):
                               {'basename': 'CVRmax', 'excl':[], 'isin':'processed', 'ext':'nii.gz', 'cmap':matplotlib.cm.jet, 'dims':(3,10)},
                           ('TMAX2STANDARD',):
                               {'basename': 'CVRdelay', 'excl':[], 'isin':'processed', 'ext':'nii.gz', 'cmap':matplotlib.cm.jet, 'dims':(3,10)},
                          }
    
reporting_folder = os.path.join(in_folder, 'reporting_images')
conversion_folder = os.path.join(reporting_folder, 'gathered')
if '5' in steps:
    ##### step 5 : reporting image generation
    ## FLAIR, CBF, CVR, CVRmax, CVRdelay
    # EtCO2 and OEF?
    
    print(f'\nStep 5: generating reporting images')
    
    thresh_names = []
    thresh_vals = []
    
    thresh_file = os.path.join(in_folder, 'thresh_vals.csv')
    
    has_thresh_file = 0
    try:
        thresh_data = pd.read_csv(thresh_file, header=None, index_col=0)
        has_thresh_file = 1
    except FileNotFoundError:
        has_ans = False
        do_search = False
        while not has_ans:
            ans = input(f'No thresh file found. Would you like to search for one?\n(y / n)\n')
            if ans == 'n':
                has_ans = True
            elif ans =='y':
                raw_id = os.path.basename(os.path.normpath(in_folder))
                split_up = raw_id.split('_')
                pt_basename = '_'.join(split_up[0:-1])
                do_search = True
                has_ans = True
            else:
                print('\nAnswer must be "y" or "n"')
                
        if do_search:
            has_ans = False
            while not has_ans:
                ans = input(f'The ID basename is {pt_basename}. Please enter a scan number (e.g., 02) that you would like to try to grab a threshhold file from, or cancel.\n(0X / cancel)\n')
                
                if ans == 'cancel':
                    print('Okay. We can make a thresh file from scratch.')
                    has_ans = True
                elif not ans.isdigit():
                    print('\nAnswer must be composed of digits only.')
                else:
                    
                    folder_to_look_for = f'{pt_basename}_{ans}'
                    print(f'Searching for folders matching {folder_to_look_for}. This may take a minute....')
                    
                    f1 = '/Users/manusdonahue/Desktop/Projects/BOLD/Data/'
                    f2 = '/Volumes/DonahueDataDrive/Data_sort/IC_Stenosis_Trial_ALL_DATA'
                    potentials = hp.find_all_folders_named(folder_to_look_for, f1)
                    potentials.extend(hp.find_all_folders_named(folder_to_look_for, f2))
                    
                    potential_threshes = [os.path.join(p, 'thresh_vals.csv') for p in potentials]
                    potential_threshes = [p for p in potential_threshes if os.path.exists(p)]
                    
                    if len(potential_threshes) == 0:
                        print("Sorry, I didn't find anything that matches.")
                    else:
                        print('I found some potential matches:')
                        for i, fi in enumerate(potential_threshes):
                            print(f'{i}:\n\t{fi}')
                        has_subans = False
                        while not has_subans:
                            subans = input('Please enter the index of the file you want to use.\n(number / cancel)\n')
                            if subans == 'cancel':
                                has_subans = True
                                continue
                            try:
                                winner = potential_threshes[int(subans)]
                                shutil.copyfile(winner, thresh_file)
                                thresh_data = pd.read_csv(thresh_file, header=None, index_col=0)
                                has_thresh_file = 1
                                has_subans = True
                                has_ans = True
                            except IndexError:
                                print('Your input must be an integer matching the indices displayed or "cancel"')
                            except ValueError:
                                print('Your input must be an integer matching the indices displayed or "cancel"')

    if os.path.exists(reporting_folder):
        shutil.rmtree(reporting_folder)
    os.mkdir(reporting_folder)
    os.mkdir(conversion_folder)

    for signature, subdict in signature_relationships.items():
        
        if has_thresh_file:
            try:
                cmax = float(thresh_data.loc[subdict['basename']])
            except KeyError:
                cmax=None
        else:
            cmax = None
            
        
        candidates = []
        # note that the signature matching includes the full path. probably not a great idea
        for subsig in signature:
            where_glob = os.path.join(in_folder, subdict['isin'], "**", f'*{subsig}*.{subdict["ext"]}')
            potential = glob.glob(where_glob, recursive=True)
            potential = [f for f in potential if not any_in_str(f, subdict['excl'])]
            candidates.extend(potential)
            
        if candidates:
            foi = candidates[-1] # pick the last in list. file of interest
        else:
            continue
        
        new_stem = f'{subdict["basename"]}.nii.gz'
        new_name = os.path.join(conversion_folder, new_stem)
        
        if subdict['ext'] == 'PAR':
            moved_name = par2nii(foi, conversion_folder)
            os.rename(moved_name, new_name)
        else:
            shutil.copy(foi, new_name)
            
        im_name = os.path.join(reporting_folder, f'{subdict["basename"]}_report_image.png')
        thresh_vals.append(nii_image(new_name, subdict['dims'], im_name, cmap=subdict['cmap'], cmax=cmax))
        thresh_names.append(subdict["basename"])
    
    thresh_dict = {key:val for key,val in zip(thresh_names, thresh_vals)}
    
    if has_thresh_file:
        try:
            thresh_dict['etco2min'] = float(thresh_data.loc['etco2min'])
            thresh_dict['etco2max'] = float(thresh_data.loc['etco2max'])
        except KeyError: # some older thresh files don't have entries for etco2
            thresh_dict['etco2min'] = 30
            thresh_dict['etco2max'] = 60
    else:
        thresh_dict['etco2min'] = 30
        thresh_dict['etco2max'] = 60
        
    thresh_ser = pd.Series(thresh_dict)
    thresh_ser.to_csv(thresh_file, header=False)
    
    
    print(f'\nReporting images generated. Elapsed time: {str_time_elapsed(start_stamp)} minutes')
    
if '6' in steps:
    ##### step 6: make the powerpoint
    
    print(f'\nStep 6: generating powerpoint')
    
    # get pt info if available
    has_age = 0
    has_dob = 0
    has_scan_date = 0
    try:
        if '.' in dobage:
            dob = dobage
            format_str = '%Y.%m.%d' # The format
            dob_dt_obj = datetime.datetime.strptime(dob, format_str)
            has_dob = 1
        else:
            pt_age = int(dobage)
            has_age = 1
    except NameError:
        pt_age = 0

    try:
        where_glob = os.path.join(in_folder, 'Acquired', "**", f'*.PAR') # just looking for any PAR
        potential = glob.glob(where_glob, recursive=True)
        read_this_one = potential[-1]
        fob = open(read_this_one)
        info = nib.parrec.parse_PAR_header(fob)
        info_dict = info[0]
        if 'exam_date' in info_dict:
            try:
                raw_scan_date = info_dict['exam_date']
                sd = raw_scan_date.split(' / ')[0]
                format_str = '%Y.%m.%d' # The format
                scan_dt_obj = datetime.datetime.strptime(sd, format_str)
                has_scan_date = 1
            except:
                print('Something went wrong with extracting the scan date, though the PAR file does seem to have a scan date')
                pt_age = 0
    except IndexError:
        where_glob = os.path.join(in_folder, 'Acquired', "**", f'*.nii*') # just looking for any NiFTI
        potential = glob.glob(where_glob, recursive=True)
        read_this_one = potential[-1]
        fob = nib.load(read_this_one)
        head = fob.header
        
        print("Unfortunately NiFTI headers do not seem to store scan dates. You'll have to set it yourself!")
        pt_age = 0
        

    if not has_scan_date:
        has_ans = False
        while not has_ans:
            ans = input(f'No scan date found. You can manually enter it now, or skip it\n(YYYY.mm.dd / skip)\n')
            if ans == 'skip':
                has_ans = True
            else:
                try:    
                    format_str = '%Y.%m.%d' # The format
                    scan_dt_obj = datetime.datetime.strptime(ans, format_str)
                    has_scan_date = 1
                    sd = ans
                    has_ans = True
                except ValueError:
                    print('\nAnswer must be "skip" or a date formatted as YYYY.mm.dd')
    

    if has_scan_date and has_dob:
        pt_age = scan_dt_obj - dob_dt_obj
        pt_age = int(pt_age.days/365.25)
        has_age = 1
        
    
    
    # make the etco2 trace
    thresh_file = os.path.join(in_folder, 'thresh_vals.csv')
    has_thresh_file = 0
    try:
        thresh_data = pd.read_csv(thresh_file, header=None, index_col=0)    
        etmin = float(thresh_data.loc['etco2min'])  
        etmax = float(thresh_data.loc['etco2max'])
        has_thresh_file = 1
    except FileNotFoundError:
        print('No thresh file found. Using default threshes for EtCO2 trace.')
        etmin = 30
        etmax = 60
    
    try:
        etco2_file = os.path.join(in_folder, 'etco2.csv')
        etco2_fig = os.path.join(reporting_folder, 'etco2.png')
        
        etco2_data = pd.read_csv(etco2_file)
        dynamics = etco2_data.iloc[:, 0]
        co2 = etco2_data.iloc[:, 1]
        
        plt.figure(figsize=((12,8)))
        plt.plot(dynamics, co2, lw=1)
        plt.scatter(dynamics, co2, color='black')
        plt.ylabel('EtCO2 (mmHg)')
        plt.xlabel('Dynamic Scan')
        plt.ylim(etmin, etmax)
        plt.tight_layout()
        plt.savefig(etco2_fig)
        plt.close()
    except FileNotFoundError:
        print(f'EtCO2 trace not found. The graph will not be generated and added to report.')
        
    
    template_loc = r'/Users/manusdonahue/Documents/Sky/repositories/scan-reporting/bin/TEMPLATE_BOLD_PLACEHOLDERS.pptx'
    template_out = os.path.join(in_folder, f'{pt_id}_report.pptx')
    
    shutil.copyfile(template_loc, template_out)
    
    replace_in_ppt('PTSTEN_###_##', pt_id, template_out)
    if has_dob:
        replace_in_ppt('dobYYYYmmdd', dob, template_out)
        print(f'DOB is {dob}')
    if has_age:
        replace_in_ppt('age##', pt_age, template_out)
        print(f'Age is {pt_age}')
    if has_scan_date:
        replace_in_ppt('scan_dateYYYYmmdd', sd, template_out)
        print(f'Scan date is {sd}')
        
    # replace_in_ppt('IMAGING', 'it worked!', template_out)
    
    #markup = os.path.join(in_folder, f'{pt_id}_report_MARKUP.pptx')
    #analyze_ppt(template_out, markup)
    
    """
    Slide 3: FLAIR
    Slide 4: CBF
    Slide 6: CVR
    Sldie 7: CVRmax
    Slide 8: CVRdelay
    Slide 10: CVR video
    Slide 11: EtCO2
    """
    
    pres = Presentation(template_out)
    
    im_names = [os.path.join(reporting_folder, f"{val['basename']}_report_image.png") for key,val in signature_relationships.items()]
    slides = [2, 3, 5, 6, 7]

    """
    movie = os.path.join(in_folder, f'{pt_id}_zstatMovie.mp4')
    movie_slide = 9
    add_ppt_image(pres.slides[movie_slide], movie, insert_type='mov', poster=im_names[0])
    """
    
    etco2_slide = 10
    
    slides.append(etco2_slide)
    im_names.append(etco2_fig)
    
    
    for slide, name in zip(slides, im_names):
        #add_ppt_image_ph(pres.slides[slide], 10, name) # don't ask why idx is 10. it for all the placeholders in this template
        try:
            add_ppt_image(pres.slides[slide], name)
        except FileNotFoundError:
            print(f'\n!!!!!\nWARNING: image {name} not found and could not be added to report\n!!!!!\n')
    
    # metrics are written as CSVs in the PSTEN_ID folder called TMAX_metrics and CBF_metrics
    # the values within are ordered as lACA, rACA, lMCA, rMCA, lPCA, rPCA
    # but the origins are MCA, ACA, PCA

    
    # plotting values by converting units to positions on a powerpoint slide
    # father forgive me for I must sin

    plot_indices = {'MCA':0, 'ACA':1, 'PCA':2}
    
    metric_names = 'lACA, rACA, lMCA, rMCA, lPCA, rPCA'.split(', ')
    plot_on = [i[1:] for i in metric_names]
    lr = [i[0] for i in metric_names]
    
    dot_sides = ['left_dot.png', 'right_dot.png']
    dot_keys = ['l', 'r']
    dot_dict = {key:os.path.join(r'/Users/manusdonahue/Documents/Sky/repositories/scan-reporting/bin', val) for key,val in zip(dot_keys, dot_sides)}
    
    file_names = ['CBF_metrics.csv', 'TMAX_metrics.csv']
    metrics_files = [os.path.join(in_folder, fn) for fn in file_names]
    slides = [4, 8]
    
    x_units_per_inch = 60 / (3.41 - 0.98) # years per inch
    y_units_per_inch_cbv = 100 / (5.5 - 2.88)
    y_units_per_inch_cvrdelay = 50 / (5.5 - 2.88)
    yupis = [y_units_per_inch_cbv, y_units_per_inch_cvrdelay]
    
    origins_cbv = [[1.03,5.51], [4.56,5.51], [8.10,5.51]] # false origins at (20yrs, 0y_units). must be adjusted
    origins_cvrd = [[0.98,5.51], [4.51,5.51], [8.05,5.51]] # false origins at (20yrs, 0y_units). must be adjusted
    adjustment = 20 / x_units_per_inch
    origins = [origins_cbv, origins_cvrd]
    for i in origins:
        for j in i:
            j[0] -= adjustment
    
    for fi, slide, yupi, origins in zip(file_names, slides, yupis, origins):
        # print(f'On slide {slide}')
        mets = pd.read_csv(os.path.join(in_folder, fi), index_col=False, header=0)
        mets = [float(i) for i in mets.columns]
        for met, plottype, side in zip(mets, plot_on, lr):
            origin = origins[plot_indices[plottype]]
            image = dot_dict[side]
            
            ad_x = pt_age / x_units_per_inch
            ad_y = met / yupi
            
            plot_dot(pres.slides[slide], image, pt_age, met,
                     origin, x_units_per_inch, yupi, size=0.11)
            
    pres.save(template_out)
    print(f'\nPowerpoint generated. Elapsed time: {str_time_elapsed(start_stamp)} minutes')

print(f'\nProcessing complete. Elapsed time: {str_time_elapsed(start_stamp)} minutes\n')